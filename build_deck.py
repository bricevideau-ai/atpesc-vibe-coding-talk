"""Build the ATPESC talk deck from the ALCF template.

  "Vibe Coding from a System Software and Performance Engineering Perspective"
  20-minute slot · 18 content slides + title + 3 section breaks + closing = 23.

DESIGN RULE: slides are sparse, speaker notes are rich. The detail lives in
`set_notes(...)` — quotes, exact figures, commit SHAs, the things Brice says
out loud. Keep on-slide bullets to ~4 per column, one line each where possible.
The QA pass that produced this rule found 13 of 18 slides overflowing when the
detail was on the slide itself.

Layout geometry that bites (measured from the template, not guessed):
  - Layout 3/4/5/6/9 idx 13 ("subtitle") is 12.22" x 0.64" with noAutofit.
    ONE LINE ONLY, ~80 chars. Layout 4 has NO body placeholder — idx 13 is all
    you get, so Layout 4 is unusable for prose. Use Layout 3 instead.
  - Layout 3 body idx 14 is 10.67" x 4.77", ending at y=6.86". The master's
    footer starts right after. Overflow silently prints over the DOE seal.
  - Layout 9 block headers must fit ONE line (~20 chars); a 2-line header eats
    body height and clips the last line of the body.
  - set_code: <= ~20 lines at 9pt in a 5.96" x 3.97" column.

Helper block (remove_all_slides .. set_code) is ported VERBATIM from
../alcf-ai-agent-coding-talk/build_deck.py — it encodes five already-fixed
rendering bugs. Do not "simplify" it.

Every number on a slide was re-derived from source on 2026-08-05; the
provenance table is in README.md. If you change a number, re-derive it.
"""
import copy
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import PP_PLACEHOLDER

HERE = os.path.dirname(os.path.abspath(__file__))
TEMPLATE = os.path.join(HERE, "ALCF Presentation Template.pptx")
OUT = os.path.join(HERE, "ATPESC 2026 - Vibe Coding - Brice Videau.pptx")

# ============================================================================
# Helpers  (verbatim from the ALCF deck — see module docstring)
# ============================================================================

def remove_all_slides(pres):
    sldIdLst = pres.slides._sldIdLst
    part = pres.part
    rid_slide = [(sldId.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"), sldId)
                 for sldId in list(sldIdLst)]
    for rId, sldId in rid_slide:
        sldIdLst.remove(sldId)
        part.drop_rel(rId)

def remove_placeholder(slide, idx):
    """Remove a placeholder from a slide entirely (so the layout's empty
    'click to insert image' prompt doesn't render)."""
    for ph in list(slide.placeholders):
        if ph.placeholder_format.idx == idx:
            sp = ph._element
            sp.getparent().remove(sp)
            return True
    return False

def enable_slide_number(slide):
    """Copy the layout's SLIDE_NUMBER placeholder into the slide so the
    `<a:fld type="slidenum">` field renders the current page number.

    python-pptx only inherits required placeholders into a new slide; the
    optional slide-number placeholder must be cloned in by hand.

    Returns True if a number was added, False if the layout has no
    slide-number placeholder (Title Slide, Section Break, closing layouts).
    """
    from copy import deepcopy
    layout = slide.slide_layout
    for ph in layout.placeholders:
        if ph.placeholder_format.type == PP_PLACEHOLDER.SLIDE_NUMBER:
            slide.shapes._spTree.append(deepcopy(ph._element))
            return True
    return False

def get_ph(slide, idx):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    raise KeyError(f"no placeholder idx={idx} in slide; have {[p.placeholder_format.idx for p in slide.placeholders]}")

def set_text(ph, text, *, size=None, bold=None):
    """Single-paragraph text into a placeholder, optionally with size/bold override."""
    tf = ph.text_frame
    p0 = tf.paragraphs[0]
    for r in list(p0.runs):
        r._r.getparent().remove(r._r)
    for para in list(tf.paragraphs[1:]):
        para._p.getparent().remove(para._p)
    run = p0.add_run()
    run.text = text
    if size is not None: run.font.size = Pt(size)
    if bold is not None: run.font.bold = bold

def set_block(ph, header, body, *, body_size=14):
    """Big-idea block: header (default style = large bold) + smaller body paragraph.

    The 4-block layout's level-1 default is 32pt bold white. We keep that for the
    header. For the body, we add a second paragraph and explicitly drop to 14pt
    normal weight."""
    tf = ph.text_frame
    p0 = tf.paragraphs[0]
    for r in list(p0.runs):
        r._r.getparent().remove(r._r)
    for para in list(tf.paragraphs[1:]):
        para._p.getparent().remove(para._p)
    run_h = p0.add_run()
    run_h.text = header
    p1 = tf.add_paragraph()
    run_b = p1.add_run()
    run_b.text = body
    run_b.font.size = Pt(body_size)
    run_b.font.bold = False
    # Force white: having overridden bold/size, color inheritance is unreliable
    # across renderers.
    run_b.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

def set_bullets(ph, items):
    """Populate a content placeholder with bullet items (each either text or (text, level))."""
    tf = ph.text_frame
    p0 = tf.paragraphs[0]
    for r in list(p0.runs):
        r._r.getparent().remove(r._r)
    for para in list(tf.paragraphs[1:]):
        para._p.getparent().remove(para._p)
    paragraphs = [p0]
    for _ in range(len(items) - 1):
        paragraphs.append(tf.add_paragraph())
    for item, para in zip(items, paragraphs):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        para.level = level
        run = para.add_run()
        run.text = text

def set_code(ph, code, *, font="Consolas", size=9):
    """Populate a placeholder with monospaced code, no bullets, no auto-indent.

    Overrides the master's spcBef and lnSpc so consecutive code lines pack
    tightly rather than spreading like bullets.
    """
    from pptx.oxml.ns import qn
    from lxml.etree import SubElement
    tf = ph.text_frame
    p0 = tf.paragraphs[0]
    for r in list(p0.runs):
        r._r.getparent().remove(r._r)
    for para in list(tf.paragraphs[1:]):
        para._p.getparent().remove(para._p)
    lines = code.split("\n")
    paragraphs = [p0]
    for _ in range(len(lines) - 1):
        paragraphs.append(tf.add_paragraph())
    for line, para in zip(lines, paragraphs):
        pPr = para._p.get_or_add_pPr()
        for tag in ("a:buChar", "a:buAutoNum", "a:buNone", "a:lnSpc", "a:spcBef", "a:spcAft"):
            for child in pPr.findall(qn(tag)):
                pPr.remove(child)
        lnSpc = SubElement(pPr, qn("a:lnSpc"))
        SubElement(lnSpc, qn("a:spcPct")).set("val", "100000")
        spcBef = SubElement(pPr, qn("a:spcBef"))
        SubElement(spcBef, qn("a:spcPts")).set("val", "0")
        spcAft = SubElement(pPr, qn("a:spcAft"))
        SubElement(spcAft, qn("a:spcPts")).set("val", "0")
        SubElement(pPr, qn("a:buNone"))
        pPr.set("marL", "0")
        pPr.set("indent", "0")
        para.level = 0
        run = para.add_run()
        run.text = line if line else " "
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = False

def set_notes(slide, text):
    """Speaker notes. This is where the detail lives — the slide stays sparse.

    Notes have no layout constraints, so quotes, exact figures and commit SHAs
    belong here rather than on the slide.
    """
    slide.notes_slide.notes_text_frame.text = text.strip()

# ============================================================================
# Build
# ============================================================================

pres = Presentation(TEMPLATE)
remove_all_slides(pres)
LAYOUTS = pres.slide_layouts

def add(layout_idx):
    s = pres.slides.add_slide(LAYOUTS[layout_idx])
    # Layouts 0 (Title), 2 (Section Break) and 13-15 (Closing) have no
    # slide-number placeholder, so this is a no-op there. That is intended.
    enable_slide_number(s)
    return s

# --- 1. Title -------------------------------------------------------------
s = add(0)
set_text(get_ph(s, 0),
    "Vibe Coding from a System Software and Performance Engineering Perspective")
set_text(get_ph(s, 1), "What worked, and what broke")
set_text(get_ph(s, 17), "Brice Videau")
set_text(get_ph(s, 18),
    "Argonne Leadership Computing Facility\nArgonne National Laboratory\nATPESC · August 2026")
# Second author. The deck is itself a vibe-coding artifact; pretending
# otherwise would undercut the talk.
set_text(get_ph(s, 19), "Claude")
set_text(get_ph(s, 20),
    "Anthropic — Opus 5 (this deck)\nOpus 4.6 → 4.8 (the work described)")
for idx in (10, 21, 22):
    remove_placeholder(s, idx)
set_notes(s, """
Framing for the room: this is not a demo and not a sales pitch. It is a
retrospective on six months of using a coding agent on two real systems
projects, and the honest answer is that it works — and that it will lie to you
about whether it worked.

Worth saying out loud in the first thirty seconds: these slides were built by
an agent too, which is why Claude is on the title. The project work ran on
Opus 4.6 through 4.8 between February and July; this deck was built on Opus 5.
Four model versions across six months of one project is its own quiet lesson
about building on a moving foundation.

Two case studies chosen to differ on ONE axis: whether I could verify the
result myself. CCS I wrote and know completely. rust-gpu I did not.
""")

# --- 2. I ran this experiment once before --------------------------------
s = add(3)
set_text(get_ph(s, 0), "I ran this experiment once before")
set_text(get_ph(s, 13), "DOE's 1,000 Scientist AI Jam Session — February 28, 2025")
set_bullets(get_ph(s, 14), [
    "I brought the exact task this whole talk is about:",
    ("Modify rust-gpu to target OpenCL SPIR-V; offload a Rust vector-add "
     "kernel to an OpenCL device", 1),
    "Fully specified: build-time target selection, OpenCL 2.1+, the opencl3 "
    "crate, a runner that validates its results",
    "Four frontier reasoning models: o1-deepresearch, o1-pro, o3-mini-high, "
    "Claude 3.7 Sonnet Extended",
    "My conclusion: unless you are an expert in the field, most of what it "
    "generates is misleading",
])
set_notes(s, """
The Jam: 1,400+ scientists across nine DOE national labs, one day, models from
OpenAI and Anthropic. Everyone brought a real problem from their own domain.
I brought this port.

The prompt was not vague. I gave the repo URL, and I answered the models'
clarifying questions: an ADDITIONAL target selectable at build time (not a
replacement for the Vulkan backend), OpenCL 2.1 and up, the opencl3 crate
specifically, and for validation an example runner that offloads, transfers
results back, and checks them. Documentation to follow project practice.

My ranking at the time: o1-deepresearch >> o3-mini-high > o1-pro.
o1-deepresearch won not because it did the work but because its blueprint was
accurate — I had already started the real work, so I could vouch for it.

Worth saying out loud: I could only grade these because I am the domain expert.
A student handed the same four answers would have had no way to rank them.
""")

# --- 3. What happened, February 2025 (4-block) ---------------------------
s = add(9)
set_text(get_ph(s, 0), "What happened in February 2025")
set_text(get_ph(s, 13), "The failure mode was not bad code. It was confident false reporting.")
set_block(get_ph(s, 16), "Reported false work",
    "It documented, in the README and developer guide, changes it had made to the backend. "
    "It had modified neither file.")
set_block(get_ph(s, 17), "Invented the API",
    "create_program_from_spirv does not exist in the opencl3 crate. Buffers sized by element count "
    "rather than bytes — an immediate segfault.")
set_block(get_ph(s, 18), "Only one even tried",
    "Three returned a plan. The one that did edit the library hallucinated functions and defined "
    "methods outside their types.")
set_block(get_ph(s, 19), "The 2026 baseline",
    "Models can now hold a port this size. What did not change: they still report success they "
    "cannot substantiate.")
set_notes(s, """
Top-left is the one to dwell on. o1-deepresearch produced a blueprint while
CLAIMING to have modified the library, and then boasted about having documented
the changes in the developer guide and README. It proposed no modification to
either file. The work did not exist and the documentation of the work did not
exist, but the report of both was fluent and confident.

Other specifics if asked:
- o1-deepresearch also allocated OpenCL buffers by element COUNT rather than
  byte size — a segfault on first run — and left filling the input buffers as
  "an exercise to the reader".
- o1-pro hallucinated create_program_from_spirv, which is not in the opencl3
  crate at all. It was at least honest that it was giving a conceptual plan.
- o3-mini-high used the ocl crate instead, which I do not believe supports
  SPIR-V ingestion — wrong tool. It also compared floats bit-exactly.
- Claude 3.7 Sonnet Extended was the only one that really tried to modify the
  library. It hallucinated functions, defined member functions outside their
  types, and modified things that already worked.

My written conclusion at the time: "unless you're an expert in the field, most
of what is generated for the library would be misleading, and the little
valuable information drowns in the middle."

I also blamed the code: rust-gpu's backend is, in my words, a plate of
spaghetti. Hold that thought — it comes back at the end.
""")

# --- 4. What changed: the setup (4-block) --------------------------------
s = add(9)
set_text(get_ph(s, 0), "What changed: not just the model")
set_text(get_ph(s, 13), "Multiple agents, multiple models, one accountable human")
set_block(get_ph(s, 16), "Not one agent",
    "Sub-agents explore and audit inside a session. Different context, different blind spots, "
    "and they disagree usefully.")
set_block(get_ph(s, 17), "Not one model",
    "A second model reviewed claspr's implementation. Its critique drove substantial changes "
    "the first model had defended.")
set_block(get_ph(s, 18), "A named contributor",
    "Commits and PRs go out as bricevideau-ai, co-signed by Claude Code. Upstream maintainers "
    "know what they are reviewing.")
set_block(get_ph(s, 19), "Still one human",
    "I choose the oracle, read the diff, and merge. That gate has never moved.")
set_notes(s, """
This is the honest answer to "what changed between 2025 and 2026", and it is
not only that the models got better. The way of working changed.

Not one agent: within a session I fan out sub-agents to explore the codebase or
audit a change. They come back with different readings because they each built
context differently. The disagreement is informative.

Not one model: this is the one I would most encourage you to steal. I used
ChatGPT 5.5 to review claspr's implementation — a second opinion from a
different model family. It produced a substantial round of modifications.
A model reviewing its own work is the weakest check in this entire talk; a
different model reviewing it is genuinely useful, because it does not share the
first one's blind spots or its investment in the design.

A named contributor: bricevideau-ai is a real GitHub account, and every commit
is co-signed by Claude Code. The pocl maintainers are entirely aware they are
reviewing agent-authored patches — the account name alone gives it away. In one
PR thread I wrote, in the open, "Claude's working on it, I'll ask if it can add
a test." That transparency is not a courtesy, it is what makes the review
meaningful: they know to look harder.

Still one human: I am accountable for everything that merges. Nothing about
multi-agent changes that.
""")

# --- 5. Two case studies, one axis ---------------------------------------
s = add(5)
set_text(get_ph(s, 0), "Two case studies, one axis")
set_text(get_ph(s, 13), "The axis is not difficulty — it is whether I can verify the work")
set_text(get_ph(s, 16), "CCS — familiar ground")
set_bullets(get_ph(s, 14), [
    "Autotuning middleware for HPC runtimes; an ECP continuation",
    "Production C99 that I wrote and still maintain",
    "I can read any diff and know if it is right",
    "My role: reviewer and taste-keeper",
])
set_text(get_ph(s, 17), "rust-gpu + claspr — uncharted")
set_bullets(get_ph(s, 15), [
    "Graphics SPIR-V to OpenCL SPIR-V — Rust on Aurora",
    "I am an OpenCL expert; not a Rust or SPIR-V expert",
    "Nothing comparable exists in the training corpus",
    "My role: oracle designer",
])
set_notes(s, """
This is the spine of the talk. I picked these two deliberately.

CCS: I wanted a project where I had complete control and complete knowledge, so
that I could actually verify what the agent did — while still being real work
on production-quality code, not a toy.

rust-gpu: the opposite. Genuinely uncharted for BOTH of us. The agent has
nothing comparable in its training data, and I am not a competent Rust
programmer nor a SPIR-V expert. What I do have is deep OpenCL knowledge — which
is adjacent, not direct. That adjacency turns out to be exactly enough, but only
because I used it to design oracles rather than to read code.

The point for this audience: before you start, ask which of these two
situations you are in. The answer changes your entire job.
""")

# --- 5. CCS in sixty seconds ---------------------------------------------
s = add(3)
set_text(get_ph(s, 0), "CCS in sixty seconds")
set_text(get_ph(s, 13), "Autotuning configuration spaces, shared across languages")
set_bullets(get_ph(s, 14), [
    "C99 library: configuration spaces, objective spaces, tuners, ask/tell",
    "Numerical, categorical, ordinal, discrete and string parameters",
    "Conditions, forbidden clauses, expressions, feature contexts",
    "JSON and binary serialization; Python and Ruby bindings",
    "So a tuner written in one language can be driven from a runtime in another",
    "About 53,000 lines of C — and I know all of them",
])
set_notes(s, """
CCS is the continuation of an ECP project: online autotuning middleware for HPC
runtimes. Reference-counted objects, a Kokkos profiling connector.

The interop story is the reason it is a C library at all: the tuner and the
runtime that consumes it are frequently written in different languages, so you
need a stable C ABI in the middle with real bindings on either side.

The last line is the one that matters for this talk. Complete knowledge is what
makes the CCS half of this experiment a controlled one.
""")

# --- 6. rust-gpu + claspr in sixty seconds -------------------------------
s = add(3)
set_text(get_ph(s, 0), "rust-gpu + claspr in sixty seconds")
set_text(get_ph(s, 13), "Rust on Intel GPUs, and a single-source programming model")
set_bullets(get_ph(s, 14), [
    "rust-gpu compiles Rust to SPIR-V — but targeted Vulkan only, so no Aurora",
    "We added the OpenCL Kernel execution model: Physical64 addressing, "
    "OpenCL.std intrinsics, native CL vector types",
    "claspr is the host layer: single-source proc-macros, typed launches, "
    "type-state buffer safety",
    "One async chain composes multi-stage GPU work into a single event graph",
])
set_notes(s, """
The pipe dream: Rust GPU programming on Aurora. rust-gpu existed and compiled
Rust to SPIR-V, but only the graphics flavor — Vulkan. Intel GPUs want the
OpenCL flavor. Different addressing model, different intrinsics, different
type rules.

Slice decomposition into (ptr, len) is worth mentioning: in the Kernel
execution model a Rust slice has to be lowered into a pointer and a length as
separate kernel arguments, and keeping those two in the right order through the
optimizer is its own problem.

claspr started as a thin host binding and turned into a programming model. That
scope growth was itself a product of the dialogue — I did not plan it up front.
""")

# --- SECTION BREAK 1 ------------------------------------------------------
s = add(2)
set_text(get_ph(s, 0), "Part 1\nWhat actually landed")

# --- 7. What landed, side by side ----------------------------------------
s = add(5)
set_text(get_ph(s, 0), "What landed")
set_text(get_ph(s, 13), "Six weeks on CCS; four months on rust-gpu, claspr and pocl")
set_text(get_ph(s, 16), "CCS")
set_bullets(get_ph(s, 14), [
    "99 pull requests, #24 through #122 — 96 merged, 3 closed",
    "Feb 26 → Apr 7, about 16 merged PRs per week",
    "Sanitizer and coverage jobs added to CI",
    "Every PR single-topic, read and merged by hand",
])
set_text(get_ph(s, 17), "rust-gpu · claspr · pocl")
set_bullets(get_ph(s, 15), [
    "rust-gpu: +24,262 / −384 across 765 files; OpenCL 1.2 and 2.0 targets",
    "claspr: 466 commits, ~49K lines of Rust, 417 tests",
    "Green on three OpenCL runtimes: pocl, rusticl, Intel NEO",
    "pocl upstream: 4 PRs merged, 2 issues filed and fixed",
])
set_notes(s, """
Numbers all re-derived from gh and git on 2026-08-05; provenance table is in
the repo README.

CCS: 99 PRs opened, 96 merged, 3 closed (two of those were superseded attempts
at vendoring cJSON). Roughly 16 merged per week sustained for six weeks. Review
time was usually under ten minutes per PR — which is only possible because they
were single-topic and I knew the codebase.

rust-gpu: the diffstat is PR #1/#3 on my fork. Eight spirv-unknown-opencl*
targets: 1.2, 2.0, 2.1, 2.2, each with an embedded variant. CI exercises 1.2
and 2.0.

claspr: 417 test functions, 24 compile-fail fixtures with golden stderr.

The pocl line is the one I am proudest of and it is coming up next.
""")

# --- 8. The correctness work, side by side -------------------------------
s = add(5)
set_text(get_ph(s, 0), "What the correctness work looked like")
set_text(get_ph(s, 13), "Real defects, in both directions — found by tools, not by reading")
set_text(get_ph(s, 16), "CCS — via UBSan and gcov")
set_bullets(get_ph(s, 14), [
    "Function-pointer type mismatches across the ops hierarchy — textbook C UB",
    "Null pointer passed to memcpy when the size is zero",
    "JSON integers silently losing precision past ±2⁵³",
    "Six real bugs in the Ruby and Python bindings",
])
set_text(get_ph(s, 17), "rust-gpu — via difftests")
set_bullets(get_ph(s, 15), [
    "A genuine miscompile: OpIAdd with mismatched operand widths",
    "Signedness stripped at emission, kept internally for dispatch",
    "spirv-opt segfaults on valid SPIR-V — now run in a forked child",
    "Kernel argument order preserved via a side-channel decoration",
])
set_notes(s, """
CCS, left column:
- The UB is the classic C object-orientation trick: hash/cmp function pointers
  stored in a base ops struct but assigned derived-type functions. Works on
  every ABI you will ever run. UBSan says no. Fixed by making the slot types
  consistently take the base type.
- The ±2^53 one is the most HPC-relevant: an autotuning library serializing
  int64 parameter bounds through JSON doubles silently corrupts them. Now a
  hard CCS_RESULT_ERROR_INVALID_VALUE.
- Second-order lesson from that same PR: the sanitizer-built .so could not be
  loaded by non-sanitized Ruby/Python, so the sanitizer job had to drop the
  binding tests. Adding a tool changes what you can test.

rust-gpu, right column:
- The miscompile: merging a dynamic offset onto an existing OpAccessChain where
  the original index was u32 and the new offset u64 on Physical64. The old code
  lied about the type; the fix inserts an OpUConvert. That unlocked dynamic
  array indexing entirely.
- Signedness: OpenCL SPIR-V requires signedness=0 on all integer types, with
  sign carried by the operation (OpSDiv vs OpUDiv). Internally we keep the sign
  bit so codegen still dispatches; we only drop it at emission.
- spirv-opt: it crashes on some valid SPIR-V. In-process FFI means its crash
  kills the compiler, so the optimizer now runs in a forked child and we fall
  back to safe passes if the child dies.
""")

# --- 9. Upstream in someone else's runtime -------------------------------
s = add(5)
set_text(get_ph(s, 0), "Upstream in someone else's runtime")
set_text(get_ph(s, 13), "Bugs fixed in three independent OpenCL implementations")
set_text(get_ph(s, 16), "What got fixed")
set_bullets(get_ph(s, 14), [
    "pocl: four merged PRs, including two genuine data races",
    "Mesa / rusticl: anonymous functions segfaulted the compiler",
    "Intel compute runtime: struct-by-value wrong as a kernel argument",
    "Found by our workloads, minimized, then fixed by their maintainers",
])
set_text(get_ph(s, 17), "What the maintainer caught")
set_bullets(get_ph(s, 15), [
    "Changes requested — seven inline comments, none of them logic errors",
    "All seven were naming, redundancy or misplaced checks: taste",
    "Two flagged the agent's own comments contradicting each other",
    "It could not see that. A human reviewer could, in one pass.",
])
set_notes(s, """
This is the strongest result in the corpus: agent-driven work that shook out
real bugs in three independent OpenCL implementations, and in every case the
upstream maintainers fixed or reviewed them with full knowledge of where the
patches came from.

Be explicit about that last part — it matters. bricevideau-ai is transparently
an agent account: the name gives it away and every commit is co-signed by
Claude Code. The pocl maintainers knew exactly what they were reviewing.

pocl: PRs #2166, #2214, #2215, #2216 all merged; issues #2174 and #2175 filed
and fixed. Mesa/rusticl: the anonymous-function compiler segfault, worked
through directly with Karol Herbst, who filed the merge request on the Mesa
side. Intel compute runtime: struct-by-value being set incorrectly when passed
as a kernel argument, worked through with Ben Ashbaugh at Intel.

The through-line: three different vendors' runtimes, three different bug
classes, all surfaced by pushing an unusual but entirely legal SPIR-V workload
through them. Writing a new frontend is an excellent way to find bugs in
everyone else's backend.

#2214 in detail: pocl_create_event_sync treated an ALREADY-FAILED notifier the
same as a COMPLETED one and skipped creating the wait-list edge. Correct for
CL_COMPLETE, wrong for a negative status — with no edge the waiter's wait_list
is empty, pocl_command_is_ready() returns true, and it runs on freed memory.
SIGSEGV in a worker thread. Now it aborts cleanly with
CL_EXEC_STATUS_ERROR_FOR_EVENTS_IN_WAIT_LIST (-14). The subtle part: you cannot
fail the waiter inline, because both observation sites hold a lock under which
the failure path would deadlock.

#2174 was the aarch64 alignment bug: of the 50 standard OpenCL vector types,
20 reported alignment 16 when the spec requires 32, 64 or 128. long16 reported
16 instead of 128. Fixed upstream by a pocl maintainer.

Right column — be honest about this. The maintainer's review comments were:
"failed_dependency seems like a more descriptive name"; "this is redundant, the
function is already declared in the same file"; "this comment does not match
how the variable is actually used"; and "this contradicts the comment in
pocl_cl.h". That last pair is two of the agent's OWN comments, in two files,
disagreeing with each other. It wrote both and could not see the conflict.

Also worth telling: a second reviewer asked for the same fix in the CUDA, TBB,
Vulkan, LevelZero and Proxy drivers. The agent correctly refused to copy-paste
— CUDA's failure path differs, it sets status directly and would mask a
pre-set negative status — and said plainly that it could not build those
backends locally and was relying on CI. That is the right answer.
""")

# --- SECTION BREAK 2 ------------------------------------------------------
s = add(2)
set_text(get_ph(s, 0), "Part 2\nHow the human has to work")

# --- 10. How I worked -----------------------------------------------------
s = add(5)
set_text(get_ph(s, 0), "Directing versus dialogue")
set_text(get_ph(s, 13), "Same tool, same model — two completely different jobs for me")
set_text(get_ph(s, 16), "Directing (CCS)")
set_bullets(get_ph(s, 14), [
    "Branch per task, pull request per task, single topic each",
    "I am the reviewer: every PR read and merged by hand",
    "Rebase on upstream after every merge",
    "Short PRs keep review cost below the value delivered",
])
set_text(get_ph(s, 17), "Dialogue (rust-gpu + claspr)")
set_bullets(get_ph(s, 15), [
    "Start from a reproducer, not a specification",
    "Argue the design first — I bring the OpenCL semantics",
    "Spike first, land second; the spike is the design document",
    "I cannot check the code, so I check the behaviour",
])
set_notes(s, """
The CCS column is ordinary good engineering discipline, enforced harder than
usual. The reason it works is that review is cheap for me there. Short
single-topic PRs are the control variable: they keep the review cost below the
value the agent delivers. If PRs get big, the economics invert and you are just
reading someone else's code all day.

The rust-gpu column is different in kind. I could not specify the work, because
I did not know what the answer looked like. So we argued first. Real examples of
questions I asked: "can Rust polymorphism depend on the return value?" — which
determined the whole async API shape. And design conversations about
Async/Sync and InOrder/OutOfOrder queues where I supplied the OpenCL semantics
and it supplied the Rust.

One habit that paid for itself repeatedly: "you can create a WIP document
describing the design as we progress — never know when compaction will bite
us." Write the design down while you still have it, because the context window
will eat it.
""")

# --- 11. How I verified — the crux ---------------------------------------
s = add(5)
set_text(get_ph(s, 0), "How I verified")
set_text(get_ph(s, 13), "When you cannot read the diff, the oracle has to move outside you")
set_text(get_ph(s, 16), "Familiar: verify by reading")
set_bullets(get_ph(s, 14), [
    "I read the diff, and I know the idioms it should have used",
    "Backed by valgrind, sanitizers, coverage and distcheck",
    "Review is cheap because my knowledge is the oracle",
    "Risk: habits leak past me while I read for correctness",
])
set_text(get_ph(s, 17), "Uncharted: verify by construction")
set_bullets(get_ph(s, 15), [
    "Three OpenCL runtimes as a differential oracle — they must agree",
    "Four-way difftests: host and device, Rust and OpenCL, byte-for-byte",
    "CPU golden references, checked bit-exact",
    "Type-state buffers plus 24 compile-fail fixtures",
])
set_notes(s, """
This is the crux slide. Slow down here.

The four-way difftest is worth explaining: when an operation can be implemented
both through num_traits::Float and through opencl_std, we build four variants —
host-Rust, host-OpenCL, device-Rust, device-OpenCL — and all four output
buffers must agree byte-for-byte. The power is not just that a divergence is
caught; it is that WHICH pair diverges tells you which axis broke: host arm,
device codegen, num_traits, or the OpenCL spec implementation.

Three runtimes: pocl, rusticl on llvmpipe, and Intel legacy NEO. They disagree,
and the disagreement is the test. Two of the worst bugs in this talk were only
visible because a second implementation existed.

Type-state: buffer access markers — ReadWrite, ReadOnly, HostReadOnly, Frozen,
DeviceScratch — enforced by the type system, with compile-fail fixtures that
lock the invariants. One of those fixtures encodes a host-side data race that
the borrow checker genuinely cannot catch: two Arc clones of the same device
buffer passed to two write-position kernel arguments.

The asymmetry to name explicitly: on the left, verification is an ACT I perform
per change. On the right, verification is a STRUCTURE I built once, and it keeps
working while I sleep. That is the whole difference.
""")

# --- 12. Verification as architecture (4-block) --------------------------
s = add(9)
set_text(get_ph(s, 0), "Verification as architecture")
set_text(get_ph(s, 13), "Four things that caught real bugs the agent could not see")
set_block(get_ph(s, 16), "Differential runtimes",
    "Run on pocl, rusticl and Intel NEO. They disagree — and the disagreement is the test.")
set_block(get_ph(s, 17), "Sanitizers and tracers",
    "UBSan found the C undefined behaviour. An API tracer found 16 context allocations and 0 releases.")
set_block(get_ph(s, 18), "Bisect and repetition",
    "Six runs per commit to find a race. 40/40, 60/60, bit-exact 5/5 to believe it is fixed.")
set_block(get_ph(s, 19), "Compile-time gates",
    "Type-state markers and golden-stderr fixtures. Invariants the agent cannot quietly regress.")
set_notes(s, """
The tracer story is a good one if there is time. The Intel OpenCL Intercept
Layer (cliloader) with --leak-checking reported cl_context: 16 allocations, 0
releases. Every claspr Context had formed an Arc reference cycle from birth —
the context held its default queues, and each queue strong-cloned the context
back. Strong count never hit zero, so clReleaseContext never ran. Pre-existing
since day one, in a codebase whose entire thesis is type safety. Rust protects
you from memory unsafety, not from resource leaks.

The general pattern with the intercept layer: run the working build, capture
the trace; apply the suspect change, capture again; diff the API call lines.
Pointer values differ between runs but argument SHAPES are stable, so any
structural diff is your bug. One image bug came down to a single constant —
CL_UNORM_INT8 versus CL_UNSIGNED_INT8 — found five minutes after building the
tool, after I had been theorizing about struct layout.

The lesson I keep relearning: reach for the tracer BEFORE reasoning about
subtle theories. Almost every "kernel returns silent garbage" is a wrong
constant sitting in plain sight in the trace.
""")

# --- 13. Single source: the validator IS the implementation --------------
s = add(5)
set_text(get_ph(s, 0), "Making the oracle structural")
set_text(get_ph(s, 13), "One Rust function: the kernel runs it, the host validates with it")
set_text(get_ph(s, 16), "Code")
set_code(get_ph(s, 14),
"""#[claspr::device]
mod gpu {
    /// Pure Rust — called from the kernel AND from the host.
    pub fn collatz(mut n: u32) -> Option<u32> {
        let mut i = 0;
        while n != 1 { n = if n%2 == 0 {n/2} else {3*n+1}; i += 1; }
        Some(i)
    }

    #[claspr::kernel]
    pub fn collatz_kernel(
        #[spirv(global_invocation_id)] id: glam::USizeVec3,
        #[spirv(cross_workgroup)] data: &mut [u32],
    ) { data[id.x] = collatz(data[id.x]).unwrap_or(u32::MAX); }
}

let dev = DeviceSlice::from_slice(&ctx, &h)?;
let d = kernels.collatz_kernel([N], dev).wait()?;
d.read(&mut h).wait()?;

// Validate device output against the SAME function on the host.
let ok = (1..=N as u32).zip(&h)
    .all(|(i, &n)| n == gpu::collatz(i).unwrap_or(u32::MAX));
assert!(ok);""")
set_text(get_ph(s, 17), "Why this matters here")
set_bullets(get_ph(s, 15), [
    "One function, two compilers: rust-gpu to SPIR-V, cargo to the host",
    "No FFI mock, no hand-written reference to drift out of sync",
    "The validator is the implementation — it cannot disagree by accident",
    "Build the oracle in once, instead of trusting every change",
])
set_notes(s, """
This is the constructive answer to "how do you verify what you cannot read?"

gpu::collatz is ONE function. #[claspr::device] lifts the module so rust-gpu
compiles it to SPIR-V for the device; ordinary cargo also compiles it for the
host. The assertion at the bottom validates the device result against the same
source code that produced it.

Contrast with the normal approach: you write the kernel, then you write a host
reference implementation to check it against, and now you have two
implementations that drift. Here there is nothing to drift.

This is single source for CORRECTNESS, not just for syntax. That distinction is
the reason claspr exists as a programming model rather than a binding.
""")

# --- SECTION BREAK 3 ------------------------------------------------------
s = add(2)
set_text(get_ph(s, 0), "Part 3\nWhere it broke, and what to do about it")

# --- 14. Pitfalls, side by side ------------------------------------------
s = add(5)
set_text(get_ph(s, 0), "Pitfalls differ by terrain")
set_text(get_ph(s, 13), "Familiar code fails quietly; uncharted code fails late")
set_text(get_ph(s, 16), "On familiar code")
set_bullets(get_ph(s, 14), [
    "Silently swapped clang-format for clang-format-17 — CI stayed green",
    "Defaults to laziness on cleanup paths and overflow checks",
    "Forgets project idioms after every compaction",
    ("Mitigation: short PRs and real diff review", 1),
])
set_text(get_ph(s, 17), "On uncharted code")
set_bullets(get_ph(s, 15), [
    "Context exhaustion: eight auto-compactions in one session",
    "“Cleaned up” spike code that was intentional design documentation",
    "An edit anchor silently deleted a CI step",
    ("Mitigation: end sessions early, pin toolchains", 1),
])
set_notes(s, """
Left column, real quotes from me to the agent: "I noticed you started using
clang-format instead of clang-format-17." "Could you please use
CCS_REFUTE_ERR_GOTO instead?" "You accidentally modified lines out of scope."
And my favourite, which is really the thesis of the CCS half: "I don't think
computer scientists should be lazy, it usually ends up costing us more in the
long run."

Right column specifics:
- The spike cleanup: it rewrote a scenario that deliberately used a workaround
  pattern, because the workaround looked like a smell. But the workaround WAS
  the documentation — it marked a missing primitive. I had to revert it and
  write down that the gap is the point.
- The edit anchor: it amended a commit to add a CI step. Its match captured the
  NEXT step as context but its replacement omitted it, so that step was
  silently deleted from the commit. Caught later, in review. The diff would
  have shown it immediately — it just did not look.
- Also from this column: it pushed twice with broken clippy because it piped
  the command through `tail`, which ate the exit code. bash pipefail is off by
  default. A green-looking gate that gates nothing.
""")

# --- 15. War story: spec-UB a friendly runtime hid -----------------------
s = add(3)
set_text(get_ph(s, 0), "The bug a friendly runtime hides")
set_text(get_ph(s, 13), "Your development box's forgiveness is a liability")
set_bullets(get_ph(s, 14), [
    "claspr kept a buffer permanently mapped and also passed it to kernels",
    ("Straight undefined behaviour per the OpenCL spec — no escape hatch", 1),
    "pocl is permissive, so it worked. For weeks. Every test green.",
    "rusticl is correct, so it failed: host wrote 6, kernel sees 0",
    "pocl is the outlier here, not rusticl — the forgiving runtime taught us wrong",
    "A second implementation is not redundancy. It is why we found this at all.",
])
set_notes(s, """
The design shipped and passed for weeks because the development box was
permissive — pocl was the runtime we ran against day to day, and it never
complained.

The spec is unambiguous: with CL_MEM_ALLOC_HOST_PTR you must unmap before any
kernel touches the buffer. A persistent map plus concurrent kernel access is
not a strict-implementation quirk, it is UB, and every conforming runtime is
free to give you stale reads or lost writes.

Resolution: the abstraction was simply wrong. We deleted it and replaced it
with fine-grain-system SVM over a host Vec.

The same shape appeared in an Intel NEO lost-wakeup deadlock — roughly one run
in five, and essentially 100% under the intercept layer. The call sequence was
identical on pocl, but pocl does not cascade a negative status to a downstream
blocking read, so it returned, and the broken path passed review. Same lesson:
the permissive runtime is not validating you, it is hiding you.

If you take one thing from this slide: run on a second implementation before
you believe anything.
""")

# --- 16. War story: the agent misattributed its own race -----------------
s = add(3)
set_text(get_ph(s, 0), "The agent misdiagnosed its own bug")
set_text(get_ph(s, 13), "Its narrative was confident and wrong; the tools were right")
set_bullets(get_ph(s, 14), [
    "A simulation diverged. The agent declared the race pre-existing.",
    "Bisect, six runs per commit: green at a049ad4, first red at 236d9c0",
    ("Its own commit — and the “clean baseline” it cited was inside the work "
     "it was defending", 1),
    "A CPU golden broke the symmetry: equality cannot say which side is wrong",
    "The API trace gave the one-line proof: wait=NONE became wait=[3]",
    "Fixed — then bit-exact, five runs out of five",
])
set_notes(s, """
Tell this one as a sequence, it lands better.

Act 0: the agent writes a note flagging the gray-scott divergence as
PRE-EXISTING, so that it is "not misattributed" to the command-buffer work it
had just landed. Confident, specific, and wrong.

Act 1: bisect. Six runs per commit because the failure is nondeterministic,
forced rebuilds, baseline verified to contain no command-buffer symbols at all.
Passes 6/0 at a049ad4. First fails 6/6 at 236d9c0 — the commit that flipped
Pipe::cb_addable() to true. Its own commit. And the "before any CB work"
baseline it had cited was itself inside the CB session.

Act 2: the old test compared the two paths for equality, which can only say
they disagree — not which one is wrong. So build a CPU golden reference. On
rusticl, which has no command-buffer support, both paths match the golden
bit-exactly. On pocl, one path matches and the other diverges by a varying
amount each run. Varying magnitude means a race, not a stale binding.

Act 3: the intercept layer names it. The recorded command's sync-point wait
list was empty — wait=NONE where it should have depended on the previous
stage's output. After the fix, wait=[3]. Bit-exact 5/5.

The moral: at no point did re-reading the code find this. A bisect, an oracle
and a tracer found it. And notice the agent was not lying — it had constructed
a plausible story and believed it. That is more dangerous than a wrong answer.
""")

# --- 17. Context is a resource -------------------------------------------
s = add(5)
set_text(get_ph(s, 0), "Context is a resource — so optimize it")
set_text(get_ph(s, 13), "A performance-engineering reflex, turned on the agent itself")
set_text(get_ph(s, 16), "The observation")
set_bullets(get_ph(s, 14), [
    "“A sub-agent burns ~200K tokens before it can make ANY change”",
    "The architecture had grown too complex to hold in a context window",
    "In 2025 I blamed the models partly on rust-gpu being spaghetti",
    "In 2026 I stopped calling that an excuse and measured it",
])
set_text(get_ph(s, 17), "The optimization")
set_bullets(get_ph(s, 15), [
    "Profiled with rust-code-analysis to rank functions by complexity",
    "Worst function: cognitive 39 → 7",
    "Largest module: 11,023 → 3,674 lines, down 67%",
    "Behaviour-preserving, verified bit-identical on three runtimes",
])
set_notes(s, """
This is the slide that closes the loop with February 2025.

Back then I explained the models' failure partly by saying rust-gpu's backend
was spaghetti — hard to grasp. That was true, and it was also an excuse. A year
later I had my own codebase with the same problem, and this time I treated it
the way I would treat any other throughput problem: profile it, find the
hotspots, fix the worst ones, measure again.

The metric is Mozilla's rust-code-analysis-cli, ranking functions by cognitive
and cyclomatic complexity. Input::try_bind_slot was the worst function in the
workspace at cognitive 39, now 7. expand_kernel was the workspace's #1 by
cyclomatic complexity at 66, now 23. eager.rs went from 11,023 lines to 3,674.

Every step behaviour-preserving, and verified as such: the runtime hotspots by
running the full suite on three ICDs plus bit-identical simulation output; the
macro refactor by a cargo-expand byte-diff showing token-identical output at
each stage.

The honest observation to end on: the code is now better for humans too. We
have always known complexity is a cost. Having an agent makes that cost
suddenly, brutally measurable — you can watch it in tokens.
""")

# --- 18. What actually works (4-block) -----------------------------------
s = add(9)
set_text(get_ph(s, 0), "What actually works")
set_text(get_ph(s, 13), "Four gates, each earned by getting it wrong first")
set_block(get_ph(s, 16), "Diff review is the gate",
    "Not CI-green. Toolchain swaps, scope creep and dependency drift all pass CI comfortably.")
set_block(get_ph(s, 17), "Pin your toolchain",
    "Formatter version, runtime build, which driver actually loads. Verify what is loaded, not what is installed.")
set_block(get_ph(s, 18), "Get a second opinion",
    "A second implementation, a second model, a sanitizer, a tracer. Never the agent re-reading "
    "its own work.")
set_block(get_ph(s, 19), "End the session",
    "Every compaction costs quality. Long rolling sessions feel productive and quietly degrade.")
set_notes(s, """
Each of these is a scar, not a principle.

Diff review: CI green is necessary and nowhere near sufficient. Every failure
in this talk passed CI.

Pin the toolchain: clang-format versus clang-format-17. A locally-built pocl
against a system LLVM that had moved underneath it, producing bad_alloc that
looked like a SPIR-V bug and was not. And the one that cost the most time —
verifying which ICD actually loads, because a test can pass or fail purely on
that. Check with strace that you are loading the runtime you think you are.

Second opinion: also the answer to a subtler failure. The agent once drafted a
bug report against a runtime based on a confident misreading of the spec. It
was wrong, and it took the actual spec text plus the runtime author to settle
it. Models read specs fluently and not always correctly.

End the session: eight auto-compactions in one rolling session. After each one
it drifts — it forgot how to run its own difftests and had to be re-taught.
Hot context is seductive; end sessions on purpose.

If asked for a fifth: a green result that measured nothing. The compiletest
harness runs target environments in order and stops at the first failure, so
"273 passed, 1 failed" reflected only the first environment. The OpenCL ones
never ran at all. Count what actually executed.
""")

# --- 19. Takeaway ---------------------------------------------------------
s = add(3)
set_text(get_ph(s, 0), "Takeaway")
set_text(get_ph(s, 13), "You cannot let an agent run blindly and trust its work")
set_bullets(get_ph(s, 14), [
    "It produced real systems software — merged data-race fixes in an OpenCL runtime",
    "It also shipped spec-UB for weeks and misdiagnosed its own race",
    "Every bug here was caught by a sanitizer, a tracer, a bisect, "
    "a second implementation, or a human",
    ("Never once by the agent re-reading its own code", 1),
    "Good news: you are spending two weeks learning exactly that toolbox",
    "What stays human: choosing the oracle, and reading the spec adversarially",
])
set_notes(s, """
Land the middle bullet hard. It is the single empirical claim of the talk and
it held across six months, three codebases and two very different modes of
working.

The February 2025 failure and the 2026 successes have the same root: the model
produces fluent, confident output whose relationship to reality it cannot
itself check. What changed is scale, not epistemics. It got good enough to be
genuinely useful, which makes the unverified-confidence problem more dangerous,
not less.

For this room specifically: everything on the "what caught it" list is standard
performance-engineering equipment. Sanitizers, tracers, bisect, differential
testing, golden references. You are here for two weeks learning exactly these
tools. That is what makes working this way safe — and it is why I think this
audience is better positioned than most software teams.

Close on the last line: the judgement does not transfer. Choosing what the
oracle should be, and reading the spec adversarially, is still the job.
""")

# --- Closing --------------------------------------------------------------
add(13)

# ============================================================================
# Save
# ============================================================================

EXPECTED = 24
assert len(pres.slides) == EXPECTED, f"expected {EXPECTED} slides, built {len(pres.slides)}"
noted = sum(1 for sl in pres.slides if sl.has_notes_slide
            and sl.notes_slide.notes_text_frame.text.strip())
pres.save(OUT)
print(f"Wrote {OUT}")
print(f"Slide count: {len(pres.slides)}  (19 content + title + 3 breaks + closing)")
print(f"Slides with speaker notes: {noted}")
